import datetime

from pptx import Presentation
from pptx.chart.chart import Plot, Chart
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LABEL_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()

# chart_data.categories = [datetime.date(2018, 10, 10), datetime.date(2018, 10, 11), datetime.date(2018, 10, 12)]
chart_data.categories = ["Catégorie 1", "Catégorie 2", "Catégorie 3", "Catégorie 4"]

serie_1 = chart_data.add_series('Série 1', (4.3, 2.5, 3.5, 4.5))
serie_2 = chart_data.add_series('Série 2', (2.4, 4.4, 1.8, 2.8))
serie_3 = chart_data.add_series('Série 3', (2, 2, 3, 5))

# grp_series = chart_data.add_series('GRP', (8, 12, 4))
# kpi_series = chart_data.add_series('KPI 1', (420, 354, 475))
# kpi2_series = chart_data.add_series('KPI 2', (320, 355, 475))
# kpi3_series = chart_data.add_series('KPI 3', (150, 25, 675))
# kpi4_series = chart_data.add_series('KPI 4', (520, 35, 275))

# grp_plot = Plot(XL_CHART_TYPE.COLUMN_CLUSTERED, [grp_series])
# kpi_plot = Plot(XL_CHART_TYPE.LINE, [kpi_series, kpi2_series, kpi3_series, kpi4_series], secondary_axis=True)

p = Plot(XL_CHART_TYPE.PIE, [])

p = Plot(XL_CHART_TYPE.COLUMN_CLUSTERED, [serie_1, serie_2, serie_3])

chart = Chart(chart_data)

# chart.add_plot(grp_plot)
# chart.add_plot(kpi_plot)
chart.add_plot(p)

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    x, y, cx, cy, chart
).chart

# chart.has_legend = False

# chart.legend.position = XL_LEGEND_POSITION.TOP

#plot = chart.plots[0]
#plot.has_data_labels = True
#data_labels = plot.data_labels

#data_labels.font.size = Pt(13)
#data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
#data_labels.position = XL_LABEL_POSITION.INSIDE_END

prs.save('chart-01.pptx')
